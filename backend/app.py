import os
from flask import Flask, request, jsonify, render_template
import google.generativeai as genai
from dotenv import load_dotenv
import PyPDF2
from pptx import Presentation

load_dotenv()

app = Flask(__name__, static_folder='static', template_folder='templates')
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
model = genai.GenerativeModel('gemini-pro')

# ---- ROUTE TO SERVE THE FRONTEND ----
@app.route('/')
def index():
    return render_template('index.html') # This serves your main page

# ---- API ENDPOINTS ----

# Endpoint 1: Process uploaded material
@app.route('/process-material', methods=['POST'])
def process_material():
    file = request.files.get('file')
    if not file:
        return jsonify({'error': 'No file uploaded'}), 400

    text = ""
    # Extract text from the uploaded file
    if file.filename.endswith('.pptx'):
        ppt = Presentation(file)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif file.filename.endswith('.pdf'):
         # NOTE: This is a basic PDF reader; more complex PDFs may need other libraries
        reader = PyPDF2.PdfReader(file.stream)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    else:
        # For .txt or other simple text files
        text = file.read().decode('utf-8')

    # Use Gemini to create study notes
    prompt = f"Summarize the following content into concise, easy-to-understand study notes with key points highlighted:\n\n{text}"
    response = model.generate_content(prompt)
    return jsonify({'summary': response.text})

# Endpoint 2: Analyze syllabus
@app.route('/analyze-syllabus', methods=['POST'])
def analyze_syllabus():
    syllabus_text = request.json['text']
    # The prompt is specifically tuned for Anna University, as requested
    prompt = f"Analyze this Anna University syllabus. Classify each topic into 'High Weightage', 'Medium Weightage', or 'Low Weightage'. Justify each classification based on typical exam patterns.\n\nSyllabus:\n{syllabus_text}"
    response = model.generate_content(prompt)
    return jsonify({'analysis': response.text})

# Endpoint 3: Generate quiz
@app.route('/generate-quiz', methods=['POST'])
def generate_quiz():
    topic_text = request.json['text']
    prompt = f"Generate a 5-question multiple-choice quiz from this text: {topic_text}. Format the output as a numbered list. For each question, provide four options (A, B, C, D) and clearly state the correct answer at the end."
    response = model.generate_content(prompt)
    return jsonify({'quiz': response.text})

if __name__ == '__main__':
    app.run(port=5000, debug=True)